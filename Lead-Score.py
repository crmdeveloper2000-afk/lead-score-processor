import logging
import requests
import io
import json
import os
import tempfile
import numpy as np
import urllib.parse
from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import Rectangle

# Configure matplotlib for server environment (no display)
matplotlib.use('Agg')

app = Flask(__name__)

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger()

# Configuration - Using environment variables for security
REFRESH_TOKEN = os.environ.get('ZOHO_REFRESH_TOKEN', "1000.97e6773099d73948b2f8e822390f5544.d500d9fbab9afb2427b86964584c4cee")
CLIENT_ID = os.environ.get('ZOHO_CLIENT_ID', "1000.8F1F1L53A19UPQPC9DTVZO9UILR2YF")
CLIENT_SECRET = os.environ.get('ZOHO_CLIENT_SECRET', "f698aef3f4c6d3c7280c3c9ecdc97b603237079321")
TEMPLATE_DOWNLOAD_URL = os.environ.get('TEMPLATE_DOWNLOAD_URL', "https://download.zoho.eu/v1/workdrive/download/llh349d6aa39cb9e347edad82405065ce5943")

@app.route('/', methods=['GET'])
def health_check():
    """
    Health check endpoint
    """
    return jsonify({
        'status': 'healthy',
        'message': 'Lead Score Processing Service is running',
        'timestamp': datetime.now().isoformat()
    }), 200

@app.route('/process-lead', methods=['POST'])
def process_lead():
    """
    Main endpoint to process lead data and generate PPT
    Expects JSON payload with lead data
    """
    try:
        # Get JSON payload from request
        if not request.is_json:
            return jsonify({'error': 'Content-Type must be application/json'}), 400
        
        payload = request.get_json()
        if not payload:
            return jsonify({'error': 'No JSON payload provided'}), 400
        
        logger.info(f"Processing payload for Lead ID: {payload.get('Lead_ID', 'Unknown')}")
        
        # Step 1: Generate access token first
        logger.info("Step 1: Generating access token...")
        access_token = get_access_token(REFRESH_TOKEN, CLIENT_ID, CLIENT_SECRET, logger)
        logger.info("Access token generated successfully!")
        
        # Step 2: Download PPT template from WorkDrive
        logger.info("Step 2: Downloading PPT template from WorkDrive...")
        ppt_template = download_ppt_template(TEMPLATE_DOWNLOAD_URL, access_token, logger)
        logger.info("PPT template downloaded successfully!")
        
        # Step 3: Replace placeholders with lead data
        logger.info("Step 3: Replacing placeholders with lead data...")
        processed_ppt = replace_placeholders(ppt_template, payload, logger)
        logger.info("Placeholders replaced successfully!")
        
        # Step 4: Add charts to slides
        logger.info("Step 4: Adding charts to slides...")
        processed_ppt = add_charts_to_slides(processed_ppt, payload, logger)
        logger.info("Charts added successfully!")
        
        # Step 5: Save PPT to temporary file for upload
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx', prefix='lead_score_')
        try:
            processed_ppt.save(temp_file.name)
            temp_path = temp_file.name
            temp_file.close()  # Ensure file is closed before upload
            
            logger.info("PPT saved to temporary file")
            
            # Step 6: Upload to Zoho WorkDrive
            logger.info("Step 6: Uploading PPT to Zoho WorkDrive...")
            upload_result = upload_to_zoho_workdrive(temp_path, access_token, logger)
            
        finally:
            # Clean up temporary file - more robust cleanup
            try:
                if 'temp_path' in locals():
                    os.unlink(temp_path)
                    logger.info("Temporary file cleaned up successfully")
            except Exception as cleanup_error:
                logger.warning(f"Could not clean up temporary file: {cleanup_error}")
        
        response_data = {
            'success': False,
            'message': '',
            'lead_id': payload.get('Lead_ID'),
            'upload_result': None,
            'attachment_result': None
        }
        
        if upload_result['success']:
            logger.info("File successfully uploaded to Zoho WorkDrive!")
            response_data['success'] = True
            response_data['upload_result'] = upload_result
            response_data['message'] = 'PPT processed and uploaded successfully'
            
            # Step 7: Attach file to Lead record if Lead_ID is available
            lead_id = payload.get('Lead_ID')
            permalink = upload_result.get('download_url')
            
            if lead_id and permalink:
                logger.info("Step 7: Attaching file to Lead record...")
                attachment_result = attach_file_to_lead(lead_id, permalink, access_token, logger)
                response_data['attachment_result'] = attachment_result
                
                if attachment_result['success']:
                    logger.info("File successfully attached to Lead!")
                    response_data['message'] += ' and attached to Lead record'
                else:
                    logger.error(f"Failed to attach file to Lead: {attachment_result['message']}")
                    response_data['message'] += ' but failed to attach to Lead record'
            else:
                logger.warning("Lead_ID or permalink not available - skipping Lead attachment")
                response_data['message'] += ' but no Lead_ID provided for attachment'
        else:
            logger.error(f"Upload failed: {upload_result['message']}")
            response_data['message'] = f"Upload failed: {upload_result['message']}"
            return jsonify(response_data), 500
        
        logger.info("Lead processing completed successfully")
        return jsonify(response_data), 200
        
    except Exception as e:
        logger.error(f"Error processing lead: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'message': 'Internal server error occurred while processing lead'
        }), 500

def get_access_token(refresh_token, client_id, client_secret, logger):
    """
    Generate access token using refresh token
    """
    try:
        # Zoho OAuth URL for EU domain
        url = "https://accounts.zoho.eu/oauth/v2/token"
        
        # Parameters for the token request
        params = {
            'refresh_token': refresh_token,
            'client_id': client_id,
            'client_secret': client_secret,
            'grant_type': 'refresh_token'
        }
        
        logger.info(f"Requesting access token from: {url}")
        
        # Make the request
        response = requests.post(url, params=params)
        
        logger.info(f"Token response status: {response.status_code}")
        logger.info(f"Token response: {response.text}")
        
        if response.status_code == 200:
            token_data = response.json()
            access_token = token_data.get('access_token')
            
            if access_token:
                logger.info("Access token generated successfully")
                return access_token
            else:
                raise Exception("No access token in response")
        else:
            raise Exception(f"Token request failed: {response.status_code} - {response.text}")
            
    except Exception as e:
        logger.error(f"Error in get_access_token: {str(e)}")
        raise

def download_ppt_template(download_url, access_token, logger):
    """
    Download PPT template from WorkDrive
    """
    try:
        # Set up headers with authorization
        headers = {
            'Authorization': f'Zoho-oauthtoken {access_token}'
        }
        
        logger.info(f"Downloading template from: {download_url}")
        
        # Download the file
        response = requests.get(download_url, headers=headers)
        
        logger.info(f"Download response status: {response.status_code}")
        
        if response.status_code == 200:
            # Load the PPT file into memory
            ppt_file = io.BytesIO(response.content)
            presentation = Presentation(ppt_file)
            
            logger.info(f"PPT loaded successfully. Slides count: {len(presentation.slides)}")
            
            return presentation
        else:
            raise Exception(f"Download failed: {response.status_code} - {response.text}")
            
    except Exception as e:
        logger.error(f"Error downloading PPT template: {str(e)}")
        raise

def upload_to_zoho_workdrive(file_path, access_token, logger):
    """
    Upload PowerPoint file to Zoho WorkDrive
    """
    try:
        # Generate meaningful filename with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        original_filename = os.path.basename(file_path)
        filename = f"lead_score_report_{timestamp}.pptx"
        
        # Zoho WorkDrive upload endpoint - EU domain to match token
        url = "https://www.zohoapis.eu/workdrive/api/v1/upload"
        
        # Headers
        headers = {
            'Authorization': f'Zoho-oauthtoken {access_token}',
            'Accept': 'application/vnd.api+json'
        }
        
        # Prepare multipart form data
        with open(file_path, 'rb') as file:
            files = {
                'content': (filename, file, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            }
            
            data = {
                'filename': filename,
                'parent_id': 'llh3437476bee74254a74bb10fa12dfe1c7ef',
                'override-name-exist': 'false'
            }
            
            logger.info(f"Uploading {filename} to Zoho WorkDrive...")
            logger.info(f"Upload URL: {url}")
            logger.info(f"Upload headers: {headers}")
            logger.info(f"Upload data: {data}")
            logger.info(f"Upload files: content=({filename}, <file_data>, {files['content'][2]})")
            
            # Upload to WorkDrive
            response = requests.post(url, headers=headers, files=files, data=data)
            
            logger.info(f"Upload response status: {response.status_code}")
            logger.info(f"Upload response: {response.text}")
            logger.info(f"Upload response headers: {dict(response.headers)}")
            
            if response.status_code == 200:  # Changed from 201 to 200
                logger.info(f"Successfully uploaded {filename} to Zoho WorkDrive")
                response_data = response.json()
                
                # Extract data from the response structure
                file_data = response_data.get('data', [{}])[0]
                attributes = file_data.get('attributes', {})
                
                return {
                    'success': True,
                    'file_id': attributes.get('resource_id'),
                    'download_url': attributes.get('Permalink'),
                    'filename': attributes.get('FileName'),
                    'message': 'File uploaded successfully'
                }
            else:
                logger.error(f"Failed to upload to WorkDrive. Status: {response.status_code}")
                logger.error(f"Response headers: {dict(response.headers)}")
                logger.error(f"Response body: {response.text}")
                return {
                    'success': False,
                    'message': f'Upload failed with status {response.status_code}: {response.text}'
                }
                
    except Exception as e:
        logger.error(f"Error uploading to Zoho WorkDrive: {str(e)}")
        return {
            'success': False,
            'message': f'Upload error: {str(e)}'
        }

def attach_file_to_lead(lead_id, permalink, access_token, logger):
    """
    Attach the uploaded PowerPoint file to the Lead record in Zoho CRM
    """
    try:
        # Don't encode the permalink URL, only encode the title if needed
        # The permalink should be used as-is
        
        # Zoho CRM attachment endpoint - EU domain
        url = f"https://crm.zoho.eu/crm/v2.1/Leads/{lead_id}/Attachments"
        
        # Headers
        headers = {
            'Authorization': f'Zoho-oauthtoken {access_token}',
            'Content-Type': 'application/json'
        }
        
        # Parameters - don't encode the attachmentUrl, only encode title if special characters
        title = 'Lead Score Matrix PPT'
        params = {
            'attachmentUrl': permalink,  # Use permalink as-is without encoding
            'title': title
        }
        
        logger.info(f"Attaching file to Lead {lead_id}...")
        logger.info(f"Attachment URL: {url}")
        logger.info(f"Permalink (unencoded): {permalink}")
        logger.info(f"Parameters: {params}")
        
        # Make the attachment request
        response = requests.post(url, headers=headers, params=params)
        
        logger.info(f"Attachment response status: {response.status_code}")
        logger.info(f"Attachment response: {response.text}")
        logger.info(f"Attachment response headers: {dict(response.headers)}")
        
        if response.status_code in [200, 201]:
            logger.info(f"Successfully attached file to Lead {lead_id}")
            response_data = response.json()
            
            return {
                'success': True,
                'attachment_id': response_data.get('data', [{}])[0].get('details', {}).get('id'),
                'message': 'File attached to Lead successfully'
            }
        else:
            logger.error(f"Failed to attach file to Lead. Status: {response.status_code}")
            logger.error(f"Response headers: {dict(response.headers)}")
            logger.error(f"Response body: {response.text}")
            return {
                'success': False,
                'message': f'Attachment failed with status {response.status_code}: {response.text}'
            }
            
    except Exception as e:
        logger.error(f"Error attaching file to Lead: {str(e)}")
        return {
            'success': False,
            'message': f'Attachment error: {str(e)}'
        }

def add_charts_to_slides(presentation, payload, logger):
    """
    Add charts to specific slides
    """
    try:
        # Add score breakdown maturity model to slide 4 (index 3)
        if len(presentation.slides) > 3:
            logger.info("Adding score breakdown maturity model to slide 4...")
            add_score_breakdown_to_slide4(presentation.slides[3], payload, logger)
        
        # Add domain scores chart to slide 5 (index 4)
        if len(presentation.slides) > 4:
            logger.info("Adding domain scores chart to slide 5...")
            add_domain_scores_chart(presentation.slides[4], payload, logger)
        
        # Add detailed reports to slide 6 (index 5)
        if len(presentation.slides) > 5:
            logger.info("Adding detailed report to slide 6...")
            add_detailed_reports_to_slide6(presentation.slides[5], payload, logger)
        
        # Add concrete recommendations report to slide 7 (index 6)
        if len(presentation.slides) > 6:
            logger.info("Adding concrete recommendations report to slide 7...")
            add_concrete_recommendations_to_slide7(presentation.slides[6], payload, logger)
        
        # Add support overview report to slide 8 (index 7)
        if len(presentation.slides) > 7:
            logger.info("Adding support overview report to slide 8...")
            add_support_overview_to_slide8(presentation.slides[7], payload, logger)
        
        return presentation
        
    except Exception as e:
        logger.error(f"Error adding charts: {str(e)}")
        raise

def add_detailed_reports_to_slide6(slide, payload, logger):
    """
    Add domain & subdomain table and spider chart side by side on slide 6
    """
    try:
        # Report 1: Domain & Subdomain table (moved 0.5 inch above and 0.5 inch left)
        logger.info("Creating domain & subdomain table...")
        subdomain_report_path = create_domain_subdomain_report(payload, logger)
        
        # Add table report - moved additional 0.5 inch to left
        left = Inches(0.1 - 0.5)  # Move 0.5 inch more left (this would be negative, so use safe minimum)
        if left.inches < 0:
            left = Inches(0.05)   # Use safe minimum left margin
        else:
            left = Inches(0.1 - 0.5)
        top = Inches(3.2 - 0.5)   # Move 0.5 inch above from current 3.2 position = 2.7
        width = Inches(9.2)       # Keep ideal size
        height = Inches(2.8)      # Keep ideal size
        
        slide.shapes.add_picture(subdomain_report_path, left, top, width, height)
        
        # Report 2: Spider chart (moved 2 inches right and 0.7 inches above for side-by-side)
        logger.info("Creating spider chart...")
        spider_chart_path = create_spider_chart_report(payload, logger)
        
        # Add spider chart moved additional 0.5 inch to right
        left = Inches(6.0 + 2.0 + 0.5)  # Move additional 0.5 inch right = 8.5
        top = Inches(3.2 - 0.7)   # Move 0.7 inches above from current 3.2 position = 2.5
        width = Inches(4.0)       # Keep width same for square aspect
        height = Inches(3.0)      # Keep current height
        
        slide.shapes.add_picture(spider_chart_path, left, top, width, height)
        
        logger.info("Both reports added side by side to slide 6")
        
    except Exception as e:
        logger.error(f"Error adding detailed reports to slide 6: {str(e)}")
        raise

def add_domain_scores_chart(slide, payload, logger):
    """
    Add domain scores table chart to slide 5
    """
    try:
        # Extract domain scores from payload
        domain_data = calculate_domain_data(payload)
        
        # Create the table chart
        chart_path = create_domain_scores_table(domain_data, logger)
        
        # Add chart to slide - centered and wider
        left = Inches(2.0)   # Move more to the left for centering
        top = Inches(2.5)    # Move up slightly
        width = Inches(9)    # Increased width
        height = Inches(4)   # Slightly increased height
        
        slide.shapes.add_picture(chart_path, left, top, width, height)
        
        # Clean up temporary file
        os.remove(chart_path)
        
        logger.info("Domain scores table added to slide 5")
        
    except Exception as e:
        logger.error(f"Error adding domain scores table: {str(e)}")
        raise

def calculate_domain_data(payload):
    """
    Calculate domain data including scores and ratings
    """
    # Map payload fields to domain data - using the actual field names from your payload
    domains_data = [
        {
            'name': 'Governance',
            'score': float(payload.get('Domain_1_Sum', 0)),
            'index': 1
        },
        {
            'name': 'Structuur',
            'score': float(payload.get('Domain_2_Sum', 0)),
            'index': 2
        },
        {
            'name': 'Proces',
            'score': float(payload.get('Domain_3_Sum', 0)),
            'index': 3
        },
        {
            'name': 'Uitkomsten & sturing',
            'score': float(payload.get('Domain_4_Sum', 0)),
            'index': 4
        }
    ]
    
    # Don't add Grand Summary row anymore
    return domains_data

def generate_star_rating(score):
    """
    Generate circle rating based on score using filled/half/empty circles
    """
    # Convert score to 0-5 scale (score/10 * 5)
    normalized_score = (score / 10) * 5
    
    # Calculate full circles
    full_circles = int(normalized_score)
    
    # Check for half circle - proper logic
    remainder = normalized_score - full_circles
    half_circle = 1 if remainder >= 0.5 else 0
    
    # Calculate empty circles
    empty_circles = 5 - full_circles - half_circle
    
    # Build circle string with reliable Unicode characters
    circle_string = '●' * full_circles  # Filled circle
    if half_circle:
        circle_string += '◐'  # Half-filled circle
    circle_string += '○' * empty_circles  # Empty circle
    
    # Add score text
    circle_string += f' ({normalized_score:.1f}/5)'
    
    return circle_string

def get_score_color(score):
    """
    Get color for score based on value (1-5 scale)
    """
    score = int(score)
    if score == 1:
        return '#dc3545'  # Red
    elif score == 2:
        return '#fd7e8a'  # Light red
    elif score == 3:
        return '#ffc107'  # Yellow/Orange (between green and red)
    elif score == 4:
        return '#90ee90'  # Light green
    elif score == 5:
        return '#28a745'  # Green
    else:
        return '#6c757d'  # Gray for unknown scores

def create_domain_subdomain_report(payload, logger):
    """
    Create beautiful domain and subdomain detailed report
    """
    try:
        # Extract individual question scores from payload
        subdomain_data = [
            {
                'domain': 'Governance',
                'subdomain': 'Visie op passende zorg',
                'score': int(payload.get('Governance_Q1_Numeric', 0)),
                'question_field': 'Governance_Q1'
            },
            {
                'domain': 'Governance', 
                'subdomain': 'Leiderschap en eigenaarschap',
                'score': int(payload.get('Governance_Q2_Numeric', 0)),
                'question_field': 'Governance_Q2'
            },
            {
                'domain': 'Structuur',
                'subdomain': 'Regionale samenwerking', 
                'score': int(payload.get('Structuur_Q1_Numeric', 0)),
                'question_field': 'Structuur_Q1'
            },
            {
                'domain': 'Structuur',
                'subdomain': 'Tools en platforms',
                'score': int(payload.get('Structuur_Q2_Numeric', 0)),
                'question_field': 'Structuur_Q2'
            },
            {
                'domain': 'Proces',
                'subdomain': 'Patiëntgericht procesontwerp',
                'score': int(payload.get('Proces_Q1_Numeric', 0)),
                'question_field': 'Proces_Q1'
            },
            {
                'domain': 'Proces',
                'subdomain': 'Leren en verbeteren',
                'score': int(payload.get('Proces_Q2_Numeric', 0)),
                'question_field': 'Proces_Q2'
            },
            {
                'domain': 'Uitkomsten & sturing',
                'subdomain': 'Outcomegericht werken',
                'score': int(payload.get('Uitkomsten_en_sturing_Q1_Numeric', 0)),
                'question_field': 'Uitkomsten_en_sturing_Q1'
            },
            {
                'domain': 'Uitkomsten & sturing',
                'subdomain': 'Monitoring en besluitvorming',
                'score': int(payload.get('Uitkomsten_en_sturing_Q2_Numeric', 0)),
                'question_field': 'Uitkomsten_en_sturing_Q2'
            }
        ]
        
        # Set up the figure with compact size for top half of slide
        fig, ax = plt.subplots(figsize=(12, 3.5))
        ax.axis('off')
        
        # Table dimensions - compact for half slide
        rows = len(subdomain_data)
        cell_height = 0.3
        cell_widths = [0.6, 3.0, 4.5, 1.2]  # Index, Domain, Subdomain, Score
        total_width = sum(cell_widths)
        
        # Start position - centered horizontally, positioned at top with no title
        start_x = (12 - total_width) / 2
        start_y = rows * cell_height + 0.2
        
        # Headers with clean styling
        headers = ['#', 'Domain', 'Subdomain', 'Score']
        header_colors = ['#2E4F99', '#3A5BA0', '#4667A7', '#5273AE']
        
        current_x = start_x
        for j, (header, width, color) in enumerate(zip(headers, cell_widths, header_colors)):
            # Clean header background
            rect = Rectangle((current_x, start_y), width, cell_height, 
                           facecolor=color, edgecolor='white', linewidth=1)
            ax.add_patch(rect)
            
            ax.text(current_x + width/2, start_y + cell_height/2, header,
                   ha='center', va='center', fontweight='bold', color='white', fontsize=11)
            current_x += width
        
        # Data rows with clean styling
        for i, data in enumerate(subdomain_data):
            current_x = start_x
            current_y = start_y - (i + 1) * cell_height
            
            # Clean alternating row colors - more white
            if i % 2 == 0:
                row_color = '#FFFFFF'  # Pure white
                border_color = '#E9ECEF'
            else:
                row_color = '#F8F9FA'  # Very light gray
                border_color = '#DEE2E6'
            
            # Index cell
            rect = Rectangle((current_x, current_y), cell_widths[0], cell_height,
                           facecolor=row_color, edgecolor=border_color, linewidth=1)
            ax.add_patch(rect)
            ax.text(current_x + cell_widths[0]/2, current_y + cell_height/2, str(i + 1),
                   ha='center', va='center', fontsize=10, fontweight='bold', color='#495057')
            current_x += cell_widths[0]
            
            # Domain cell with domain-specific colors
            domain_colors = {
                'Governance': '#6F42C1',
                'Structuur': '#20C997', 
                'Proces': '#FD7E14',
                'Uitkomsten & sturing': '#DC3545'
            }
            domain_color = domain_colors.get(data['domain'], '#6C757D')
            
            rect = Rectangle((current_x, current_y), cell_widths[1], cell_height,
                           facecolor=row_color, edgecolor=border_color, linewidth=1)
            ax.add_patch(rect)
            ax.text(current_x + cell_widths[1]/2, current_y + cell_height/2, data['domain'],
                   ha='center', va='center', fontsize=10, fontweight='bold', color=domain_color)
            current_x += cell_widths[1]
            
            # Subdomain cell
            rect = Rectangle((current_x, current_y), cell_widths[2], cell_height,
                           facecolor=row_color, edgecolor=border_color, linewidth=1)
            ax.add_patch(rect)
            ax.text(current_x + cell_widths[2]/2, current_y + cell_height/2, data['subdomain'],
                   ha='center', va='center', fontsize=9, color='#212529')
            current_x += cell_widths[2]
            
            # Score cell with clean score display
            score_color = get_score_color(data['score'])
            rect = Rectangle((current_x, current_y), cell_widths[3], cell_height,
                           facecolor=score_color, edgecolor='white', linewidth=1)
            ax.add_patch(rect)
            
            ax.text(current_x + cell_widths[3]/2, current_y + cell_height/2, str(data['score']),
                   ha='center', va='center', fontsize=12, fontweight='bold', color='white')
        
        # Set limits for compact display without title
        ax.set_xlim(0, 12)
        ax.set_ylim(0, start_y + 0.3)
        
        # Save to temporary file with transparent background
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                   facecolor='none', edgecolor='none', pad_inches=0.05, transparent=True)
        plt.close()
        
        logger.info(f"Domain & Subdomain report created: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logger.error(f"Error creating domain & subdomain report: {str(e)}")
        raise

def create_spider_chart_report(payload, logger):
    """
    Create beautiful compact spider/radar chart for subdomain performance analysis
    """
    try:
        # Extract subdomain scores from payload (same as table report)
        subdomain_data = [
            {
                'subdomain': 'Visie op passende zorg',
                'score': int(payload.get('Governance_Q1_Numeric', 0))
            },
            {
                'subdomain': 'Leiderschap en eigenaarschap',
                'score': int(payload.get('Governance_Q2_Numeric', 0))
            },
            {
                'subdomain': 'Regionale samenwerking', 
                'score': int(payload.get('Structuur_Q1_Numeric', 0))
            },
            {
                'subdomain': 'Tools en platforms',
                'score': int(payload.get('Structuur_Q2_Numeric', 0))
            },
            {
                'subdomain': 'Patiëntgericht procesontwerp',
                'score': int(payload.get('Proces_Q1_Numeric', 0))
            },
            {
                'subdomain': 'Leren en verbeteren',
                'score': int(payload.get('Proces_Q2_Numeric', 0))
            },
            {
                'subdomain': 'Sturing op uitkomsten',
                'score': int(payload.get('Uitkomsten_en_sturing_Q1_Numeric', 0))
            },
            {
                'subdomain': 'Datagebruik in besluitvorming',
                'score': int(payload.get('Uitkomsten_en_sturing_Q2_Numeric', 0))
            }
        ]
        
        # Set up the figure - simple and clean
        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(projection='polar'))
        fig.patch.set_facecolor('white')
        fig.patch.set_alpha(0)
        
        # Extract data for radar chart
        categories = [item['subdomain'] for item in subdomain_data]
        scores = [item['score'] for item in subdomain_data]
        
        # Create clean, readable labels
        clean_labels = [
            "Visie op passende zorg",
            "Leiderschap en eigenaarschap", 
            "Regionale samenwerking",
            "Tools en platforms",
            "Patiëntgericht procesontwerp",
            "Leren en verbeteren",
            "Sturing op uitkomsten",
            "Datagebruik in besluitvorming"
        ]
        
        # Number of variables
        N = len(categories)
        
        # Compute angle for each axis
        angles = [n / float(N) * 2 * np.pi for n in range(N)]
        angles += angles[:1]  # Complete the circle
        
        # Add scores and complete the circle
        scores += scores[:1]
        
        # Set up clean polar grid
        ax.set_ylim(0, 5)
        ax.set_rticks([])  # Remove the radial tick numbers 
        ax.grid(True, alpha=0.3, linestyle='-', linewidth=0.5, color='#CCCCCC')
        ax.set_facecolor('white')
        
        # Different colors for each border segment - minimal width
        colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DDA0DD', '#98D8C8', '#F7DC6F']
        
        # Plot each segment with different colored borders - minimal width
        for i in range(N):
            start_angle = angles[i]
            end_angle = angles[i+1]
            start_score = scores[i] 
            end_score = scores[i+1]
            
            # Draw colored border for each segment
            ax.plot([start_angle, end_angle], [start_score, end_score], 
                   color=colors[i % len(colors)], linewidth=1.5)
        
        # Fill the area with light blue
        ax.fill(angles, scores, alpha=0.2, color='#4A90E2')
        
        # Add score values at each point - blue numbers
        for angle, score in zip(angles[:-1], scores[:-1]):
            if score > 0:
                ax.text(angle, score + 0.2, f'{score}', ha='center', va='center', 
                       fontsize=10, fontweight='bold', color='#4A90E2')
        
        # Set up category labels - positioned normally
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(clean_labels, fontsize=9, color='#333333', fontweight='normal')
        
        # Clean outer circle
        ax.spines['polar'].set_visible(True)
        ax.spines['polar'].set_linewidth(1)
        ax.spines['polar'].set_color('#CCCCCC')
        
        # Adjust layout to prevent label cutoff
        plt.tight_layout()
        
        # Save to temporary file with transparent background and smaller size
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                   facecolor='none', edgecolor='none', transparent=True, pad_inches=0.1)
        plt.close()
        
        logger.info(f"Beautiful compact spider chart created: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logger.error(f"Error creating spider chart report: {str(e)}")
        raise

def add_concrete_recommendations_to_slide7(slide, payload, logger):
    """
    Add concrete recommendations report to slide 7
    """
    try:
        # Create concrete recommendations report
        recommendations_chart_path = create_concrete_recommendations_report(payload, logger)
        
        # Position the chart to cover most of the slide
        left = Inches(1.0)  # 1.0 inch from left
        top = Inches(2.7)   # 2.7 inches from top (moved 0.5 inch up from 3.2)
        width = Inches(10.5)  # 10.5 inches wide (reduced by 2 inches)
        height = Inches(3.6)  # 3.6 inches high (increased by 0.6 inches from 3.0)
        
        # Add the chart to slide
        slide.shapes.add_picture(recommendations_chart_path, left, top, width, height)
        
        logger.info("Concrete recommendations report added to slide 7")
        
    except Exception as e:
        logger.error(f"Error adding concrete recommendations to slide 7: {str(e)}")
        raise

def create_concrete_recommendations_report(payload, logger):
    """
    Create concrete recommendations table based on scores <= 3
    """
    try:
        # Extract subdomain scores from payload
        subdomain_scores = {}
        
        # Map payload fields to subdomain names (based on your query logic)
        field_mapping = {
            "Governance_Q1": "Visie op passende zorg",
            "Governance_Q2": "Leiderschap en eigenaarschap", 
            "Structuur_Q1": "Regionale samenwerking",
            "Structuur_Q2": "Tools en platforms",
            "Proces_Q1": "Patiëntgericht procesontwerp",
            "Proces_Q2": "Leren en verbeteren",
            "Uitkomsten_en_sturing_Q1": "Outcomegericht werken",
            "Uitkomsten_en_sturing_Q2": "Monitoring en besluitvorming"
        }
        
        # Convert text scores to numeric values
        for field, subdomain in field_mapping.items():
            if field in payload:
                score_text = payload[field]
                # Extract number from text like "1. Some description"
                if isinstance(score_text, str) and score_text.strip():
                    score = int(score_text.split('.')[0])
                    subdomain_scores[subdomain] = score
                elif field.endswith('_Numeric') and field.replace('_Numeric', '') in field_mapping:
                    # Handle numeric fields
                    score = int(payload[field])
                    base_field = field.replace('_Numeric', '')
                    if base_field in field_mapping:
                        subdomain_scores[field_mapping[base_field]] = score
        
        # Check for numeric fields in payload
        numeric_fields = {
            "Structuur_Q1_Numeric": "Regionale samenwerking",
            "Structuur_Q2_Numeric": "Tools en platforms", 
            "Uitkomsten_en_sturing_Q1_Numeric": "Outcomegericht werken",
            "Uitkomsten_en_sturing_Q2_Numeric": "Monitoring en besluitvorming"
        }
        
        for field, subdomain in numeric_fields.items():
            if field in payload:
                score = int(payload[field])
                subdomain_scores[subdomain] = score
        
        # Define advice mapping based on your HTML attachment
        advice_mapping = {
            ("Visie op passende zorg", 1): {
                "advice": "Faciliteer een visie- en inspiratiesessie met stakeholders.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Visie op passende zorg", 2): {
                "advice": "Vertaal losse ideeën naar een eerste conceptvisie.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Visie op passende zorg", 3): {
                "advice": "Verscherp visie en koppel concrete doelen en termijnen.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Leiderschap en eigenaarschap", 1): {
                "advice": "Benoem een bestuurlijk ambassadeur voor passende zorg.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Leiderschap en eigenaarschap", 2): {
                "advice": "Betrek bestuur actief bij voortgang en beslismomenten.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Leiderschap en eigenaarschap", 3): {
                "advice": "Geef bestuur formele rol in governance.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Regionale samenwerking", 1): {
                "advice": "Breng partners in kaart en start eerste verkenningsgesprekken.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Regionale samenwerking", 2): {
                "advice": "Organiseer maandelijks thematisch netwerkoverleg.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Regionale samenwerking", 3): {
                "advice": "Versterk met gezamenlijke doelen en actielijst.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Tools en platforms", 1): {
                "advice": "Start met gedeelde mappen of formats.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Tools en platforms", 2): {
                "advice": "Verken dashboard-tools of Zoho/PowerBI.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Tools en platforms", 3): {
                "advice": "Versnel ontwikkeling en test actief met gebruikers.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Patiëntgericht procesontwerp", 1): {
                "advice": "Visualiseer de patiëntreis met team of patiëntpanel.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Patiëntgericht procesontwerp", 2): {
                "advice": "Evalueer pilots en werk verbeterideeën verder uit.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Patiëntgericht procesontwerp", 3): {
                "advice": "Standaardiseer het proces en monitor op resultaat.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Leren en verbeteren", 1): {
                "advice": "Start met reflectie- of verbetermomenten per kwartaal.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Leren en verbeteren", 2): {
                "advice": "Implementeer eenvoudige PDCA-cyclus op teamniveau.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Leren en verbeteren", 3): {
                "advice": "Borg deze in overleggen en dashboards.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Outcomegericht werken", 1): {
                "advice": "Definieer 2–3 relevante uitkomstindicatoren.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Outcomegericht werken", 2): {
                "advice": "Maak ze zichtbaar in teamoverleg of dashboard.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Outcomegericht werken", 3): {
                "advice": "Koppel outcome aan proces- en beslisinformatie.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            },
            ("Monitoring en besluitvorming", 1): {
                "advice": "Start met een maandelijks stuurmoment.",
                "support": "Startsessie of training: visie, netwerk of dashboard opzetten.",
                "support_type": "Training"
            },
            ("Monitoring en besluitvorming", 2): {
                "advice": "Introduceer KPI-dashboard met kwartaalupdate.",
                "support": "Co-creatie workshop: structuur of pilotplan uitwerken.",
                "support_type": "Workshop"
            },
            ("Monitoring en besluitvorming", 3): {
                "advice": "Train teams in gebruik en interpretatie.",
                "support": "Consultancy: concretiseer aanpak en borg werkwijze.",
                "support_type": "Consultancy"
            }
        }
        
        # Filter recommendations for scores <= 3
        recommendations = []
        organization = payload.get("Organization", "Organization")
        first_name = payload.get("First_Name", "")
        last_name = payload.get("Last_Name", "")
        full_name = f"{first_name} {last_name}".strip()
        
        # Add domain mapping
        domain_mapping = {
            "Visie op passende zorg": "Governance",
            "Leiderschap en eigenaarschap": "Governance",
            "Regionale samenwerking": "Structuur", 
            "Tools en platforms": "Structuur",
            "Patiëntgericht procesontwerp": "Proces",
            "Leren en verbeteren": "Proces",
            "Outcomegericht werken": "Uitkomsten & sturing",
            "Monitoring en besluitvorming": "Uitkomsten & sturing"
        }
        
        for subdomain, score in subdomain_scores.items():
            if score <= 3 and (subdomain, score) in advice_mapping:
                rec = advice_mapping[(subdomain, score)]
                recommendations.append({
                    "full_name": full_name,
                    "organization": organization,
                    "subdomain": subdomain,
                    "domain": domain_mapping.get(subdomain, ""),
                    "score": score,
                    "advice": rec["advice"],
                    "support": rec["support"],
                    "support_type": rec["support_type"]
                })
        
        if not recommendations:
            # If no recommendations (all scores > 3), create a congratulatory message
            fig, ax = plt.subplots(figsize=(14, 8))
            ax.axis('off')
            
            ax.text(0.5, 0.5, "Gefeliciteerd! Alle scores zijn hoger dan 3.\nGeen concrete aanbevelingen nodig.", 
                   ha='center', va='center', fontsize=20, fontweight='bold', color='#2E8B57')
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none', transparent=False)
            plt.close()
            
            logger.info(f"No recommendations chart created: {temp_file.name}")
            return temp_file.name
        
        # Create the table
        fig, ax = plt.subplots(figsize=(14, max(3, len(recommendations) * 0.4 + 1)))
        ax.axis('off')
        
        # Define table structure
        headers = ["Domain", "Subdomain", "Advice", "Score"]
        col_widths = [0.20, 0.30, 0.40, 0.10]  # Proportional column widths
        
        # Color scheme - darker and more visible
        header_color = '#2C5282'  # Darker blue
        row_colors = ['#E2E8F0', '#F7FAFC']  # Darker alternating row colors
        score_colors = {1: '#DC3545', 2: '#FF6B6B', 3: '#FF8C00'}  # Red, Light Red, Orange
        
        # Calculate positions
        table_width = 0.95
        table_left = 0.025
        base_row_height = 0.08
        header_height = 0.10
        
        # Pre-calculate row heights based on content
        row_heights = []
        for rec in recommendations:
            max_lines = 1
            # Check advice column for text wrapping
            advice_text = rec["advice"]
            if len(advice_text) > 45:  # Shorter threshold for advice column
                words = advice_text.split(' ')
                current_line = ""
                line_count = 0
                for word in words:
                    test_line = current_line + " " + word if current_line else word
                    if len(test_line) <= 45:
                        current_line = test_line
                    else:
                        if current_line:
                            line_count += 1
                        current_line = word
                if current_line:
                    line_count += 1
                max_lines = max(max_lines, line_count)
            
            # Calculate row height based on content
            row_height = max(base_row_height, base_row_height * max_lines * 0.9)
            row_heights.append(row_height)
        
        # Draw headers - no title needed
        header_y = 0.90  # Start headers higher since no title
        current_x = table_left
        
        for i, (header, width) in enumerate(zip(headers, col_widths)):
            # Header background
            rect = plt.Rectangle((current_x, header_y - header_height/2), 
                               width * table_width, header_height,
                               facecolor=header_color, alpha=1.0, edgecolor='#1A365D', linewidth=1)  # Solid header with darker border
            ax.add_patch(rect)
            
            # Header text
            ax.text(current_x + (width * table_width)/2, header_y, header,
                   ha='center', va='center', fontweight='bold', color='white', fontsize=11)
            
            current_x += width * table_width
        
        # Draw data rows with dynamic heights
        current_y = header_y - header_height/2
        for row_idx, rec in enumerate(recommendations):
            row_height = row_heights[row_idx]
            row_y = current_y - row_height/2
            current_x = table_left
            
            # Alternate row colors
            row_color = row_colors[row_idx % 2]
            
            # Row data
            row_data = [
                rec["domain"],
                rec["subdomain"],
                rec["advice"],
                str(rec["score"])
            ]
            
            for col_idx, (data, width) in enumerate(zip(row_data, col_widths)):
                # Cell background
                cell_color = row_color
                if col_idx == 3:  # Score column (now at position 3)
                    cell_color = score_colors.get(rec["score"], row_color)
                
                rect = plt.Rectangle((current_x, row_y - row_height/2), 
                                   width * table_width, row_height,
                                   facecolor=cell_color, alpha=0.9, edgecolor='#CBD5E0', linewidth=1)  # Darker borders and higher alpha
                ax.add_patch(rect)
                
                # Cell text
                text_color = 'white' if col_idx == 3 else '#1A202C'  # Darker text for better contrast
                font_weight = 'bold' if col_idx == 3 else 'normal'
                font_size = 9
                
                # Wrap text for advice column only
                if col_idx == 2 and len(data) > 45:
                    wrapped_lines = []
                    words = data.split(' ')
                    current_line = ""
                    for word in words:
                        test_line = current_line + " " + word if current_line else word
                        if len(test_line) <= 45:
                            current_line = test_line
                        else:
                            if current_line:
                                wrapped_lines.append(current_line)
                            current_line = word
                    if current_line:
                        wrapped_lines.append(current_line)
                    
                    # Display wrapped text with dynamic spacing
                    line_spacing = row_height / max(len(wrapped_lines), 1) * 0.6
                    start_y = row_y + (len(wrapped_lines) - 1) * line_spacing / 2
                    for line_idx, line in enumerate(wrapped_lines):
                        ax.text(current_x + (width * table_width)/2, 
                               start_y - line_idx * line_spacing, line,
                               ha='center', va='center', fontweight=font_weight, 
                               color=text_color, fontsize=font_size-1)
                else:
                    ax.text(current_x + (width * table_width)/2, row_y, data,
                           ha='center', va='center', fontweight=font_weight, 
                           color=text_color, fontsize=font_size)
                
                current_x += width * table_width
            
            # Update current_y for next row
            current_y -= row_height
        
        # Set limits and remove axes
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        
        plt.tight_layout()
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                   facecolor='white', edgecolor='none', transparent=False, pad_inches=0.1)
        plt.close()
        
        logger.info(f"Concrete recommendations report created with {len(recommendations)} recommendations: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logger.error(f"Error creating concrete recommendations report: {str(e)}")
        raise

def add_support_overview_to_slide8(slide, payload, logger):
    """
    Add support overview report to slide 8
    """
    try:
        # Create support overview report
        support_overview_chart_path = create_support_overview_report(payload, logger)
        
        # Position the chart same as slide 7
        left = Inches(1.0)  # 1.0 inch from left
        top = Inches(2.7)   # 2.7 inches from top
        width = Inches(10.5)  # 10.5 inches wide
        height = Inches(3.6)  # 3.6 inches high
        
        # Add the chart to slide
        slide.shapes.add_picture(support_overview_chart_path, left, top, width, height)
        
        logger.info("Support overview report added to slide 8")
        
    except Exception as e:
        logger.error(f"Error adding support overview to slide 8: {str(e)}")
        raise

def create_support_overview_report(payload, logger):
    """
    Create support overview table based on scores <= 3
    """
    try:
        # Extract subdomain scores from payload
        subdomain_scores = {}
        
        # Map payload fields to subdomain names
        field_mapping = {
            "Governance_Q1": "Visie op passende zorg",
            "Governance_Q2": "Leiderschap en eigenaarschap", 
            "Structuur_Q1": "Regionale samenwerking",
            "Structuur_Q2": "Tools en platforms",
            "Proces_Q1": "Patiëntgericht procesontwerp",
            "Proces_Q2": "Leren en verbeteren",
            "Uitkomsten_en_sturing_Q1": "Outcomegericht werken",
            "Uitkomsten_en_sturing_Q2": "Monitoring en besluitvorming"
        }
        
        # Convert text scores to numeric values
        for field, subdomain in field_mapping.items():
            if field in payload:
                score_text = payload[field]
                # Extract number from text like "1. Some description"
                if isinstance(score_text, str) and score_text.strip():
                    score = int(score_text.split('.')[0])
                    subdomain_scores[subdomain] = score
        
        # Check for numeric fields in payload
        numeric_fields = {
            "Governance_Q1_Numeric": "Visie op passende zorg",
            "Governance_Q2_Numeric": "Leiderschap en eigenaarschap",
            "Structuur_Q1_Numeric": "Regionale samenwerking",
            "Structuur_Q2_Numeric": "Tools en platforms", 
            "Proces_Q1_Numeric": "Patiëntgericht procesontwerp",
            "Proces_Q2_Numeric": "Leren en verbeteren",
            "Uitkomsten_en_sturing_Q1_Numeric": "Outcomegericht werken",
            "Uitkomsten_en_sturing_Q2_Numeric": "Monitoring en besluitvorming"
        }
        
        for field, subdomain in numeric_fields.items():
            if field in payload:
                score = int(payload[field])
                subdomain_scores[subdomain] = score
        
        # Define support mapping based on your HTML attachment
        support_mapping = {
            ("Visie op passende zorg", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Visie op passende zorg", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Visie op passende zorg", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Leiderschap en eigenaarschap", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Leiderschap en eigenaarschap", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Leiderschap en eigenaarschap", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Regionale samenwerking", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Regionale samenwerking", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Regionale samenwerking", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Tools en platforms", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Tools en platforms", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Tools en platforms", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Patiëntgericht procesontwerp", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Patiëntgericht procesontwerp", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Patiëntgericht procesontwerp", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Leren en verbeteren", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Leren en verbeteren", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Leren en verbeteren", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Outcomegericht werken", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Outcomegericht werken", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Outcomegericht werken", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            },
            ("Monitoring en besluitvorming", 1): {
                "support_type": "Training",
                "description": "Startsessie of training: visie, netwerk of dashboard opzetten."
            },
            ("Monitoring en besluitvorming", 2): {
                "support_type": "Workshop",
                "description": "Co-creatie workshop: structuur of pilotplan uitwerken."
            },
            ("Monitoring en besluitvorming", 3): {
                "support_type": "Consultancy",
                "description": "Consultancy: concretiseer aanpak en borg werkwijze."
            }
        }
        
        # Filter support recommendations for scores <= 3
        support_recommendations = []
        
        for subdomain, score in subdomain_scores.items():
            if score <= 3 and (subdomain, score) in support_mapping:
                support_info = support_mapping[(subdomain, score)]
                support_recommendations.append({
                    "subdomain": subdomain,
                    "support_type": support_info["support_type"],
                    "description": support_info["description"],
                    "score": score
                })
        
        if not support_recommendations:
            # If no support needed (all scores > 3), create a congratulatory message
            fig, ax = plt.subplots(figsize=(14, 8))
            ax.axis('off')
            
            ax.text(0.5, 0.5, "Gefeliciteerd! Alle scores zijn hoger dan 3.\nGeen ondersteuning nodig.", 
                   ha='center', va='center', fontsize=20, fontweight='bold', color='#2E8B57')
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none', transparent=False)
            plt.close()
            
            logger.info(f"No support needed chart created: {temp_file.name}")
            return temp_file.name
        
        # Create the table
        fig, ax = plt.subplots(figsize=(14, max(3, len(support_recommendations) * 0.4 + 1)))
        ax.axis('off')
        
        # Define table structure
        headers = ["Subdomain", "Support Type", "Description", "Score"]
        col_widths = [0.25, 0.15, 0.50, 0.10]  # Proportional column widths
        
        # Color scheme - darker and more visible
        header_color = '#2C5282'  # Darker blue
        row_colors = ['#E2E8F0', '#F7FAFC']  # Darker alternating row colors
        score_colors = {1: '#DC3545', 2: '#FF6B6B', 3: '#FF8C00'}  # Dark Red, Light Red, Orange
        
        # Calculate positions
        table_width = 0.95
        table_left = 0.025
        base_row_height = 0.08
        header_height = 0.10
        
        # Pre-calculate row heights based on content
        row_heights = []
        for rec in support_recommendations:
            max_lines = 1
            # Check description column for text wrapping
            description_text = rec["description"]
            if len(description_text) > 60:
                words = description_text.split(' ')
                current_line = ""
                line_count = 0
                for word in words:
                    test_line = current_line + " " + word if current_line else word
                    if len(test_line) <= 60:
                        current_line = test_line
                    else:
                        if current_line:
                            line_count += 1
                        current_line = word
                if current_line:
                    line_count += 1
                max_lines = max(max_lines, line_count)
            
            # Calculate row height based on content
            row_height = max(base_row_height, base_row_height * max_lines * 0.9)
            row_heights.append(row_height)
        
        # Draw headers - no title needed
        header_y = 0.90  # Start headers higher since no title
        current_x = table_left
        
        for i, (header, width) in enumerate(zip(headers, col_widths)):
            # Header background
            rect = plt.Rectangle((current_x, header_y - header_height/2), 
                               width * table_width, header_height,
                               facecolor=header_color, alpha=1.0, edgecolor='#1A365D', linewidth=1)  # Solid header with darker border
            ax.add_patch(rect)
            
            # Header text
            ax.text(current_x + (width * table_width)/2, header_y, header,
                   ha='center', va='center', fontweight='bold', color='white', fontsize=11)
            
            current_x += width * table_width
        
        # Draw data rows with dynamic heights
        current_y = header_y - header_height/2
        for row_idx, rec in enumerate(support_recommendations):
            row_height = row_heights[row_idx]
            row_y = current_y - row_height/2
            current_x = table_left
            
            # Alternate row colors
            row_color = row_colors[row_idx % 2]
            
            # Row data
            row_data = [
                rec["subdomain"],
                rec["support_type"],
                rec["description"],
                str(rec["score"])
            ]
            
            for col_idx, (data, width) in enumerate(zip(row_data, col_widths)):
                # Cell background
                cell_color = row_color
                if col_idx == 3:  # Score column (at position 3)
                    cell_color = score_colors.get(rec["score"], row_color)
                
                rect = plt.Rectangle((current_x, row_y - row_height/2), 
                                   width * table_width, row_height,
                                   facecolor=cell_color, alpha=0.9, edgecolor='#CBD5E0', linewidth=1)  # Darker borders and higher alpha
                ax.add_patch(rect)
                
                # Cell text
                text_color = 'white' if col_idx == 3 else '#1A202C'  # Darker text for better contrast
                font_weight = 'bold' if col_idx == 3 else 'normal'
                font_size = 9
                
                # Wrap text for description column only
                if col_idx == 2 and len(data) > 60:
                    wrapped_lines = []
                    words = data.split(' ')
                    current_line = ""
                    for word in words:
                        test_line = current_line + " " + word if current_line else word
                        if len(test_line) <= 60:
                            current_line = test_line
                        else:
                            if current_line:
                                wrapped_lines.append(current_line)
                            current_line = word
                    if current_line:
                        wrapped_lines.append(current_line)
                    
                    # Display wrapped text with dynamic spacing
                    line_spacing = row_height / max(len(wrapped_lines), 1) * 0.6
                    start_y = row_y + (len(wrapped_lines) - 1) * line_spacing / 2
                    for line_idx, line in enumerate(wrapped_lines):
                        ax.text(current_x + (width * table_width)/2, 
                               start_y - line_idx * line_spacing, line,
                               ha='center', va='center', fontweight=font_weight, 
                               color=text_color, fontsize=font_size-1)
                else:
                    ax.text(current_x + (width * table_width)/2, row_y, data,
                           ha='center', va='center', fontweight=font_weight, 
                           color=text_color, fontsize=font_size)
                
                current_x += width * table_width
            
            # Update current_y for next row
            current_y -= row_height
        
        # Set limits and remove axes
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        
        plt.tight_layout()
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                   facecolor='white', edgecolor='none', transparent=False, pad_inches=0.1)
        plt.close()
        
        logger.info(f"Support overview report created with {len(support_recommendations)} recommendations: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logger.error(f"Error creating support overview report: {str(e)}")
        raise

def add_score_breakdown_to_slide4(slide, payload, logger):
    """
    Add score breakdown maturity model visualization to slide 4
    """
    try:
        # Create score breakdown chart
        score_breakdown_chart_path = create_score_breakdown_chart(payload, logger)
        
        # Position the chart to cover most of the slide
        left = Inches(0.5)    # Small margin from left
        top = Inches(1.5 - 0.9)     # Move up 0.9 inch = 0.6 from top
        width = Inches(12.0)  # Wide chart
        height = Inches(5.5 - 1.5)  # Reduce height by 1.5 inches = 4.0
        
        # Add the chart to slide
        slide.shapes.add_picture(score_breakdown_chart_path, left, top, width, height)
        
        logger.info("Score breakdown maturity model added to slide 4")
        
    except Exception as e:
        logger.error(f"Error adding score breakdown to slide 4: {str(e)}")
        raise

def create_score_breakdown_chart(payload, logger):
    """
    Create comprehensive score breakdown maturity model visualization
    """
    try:
        # Calculate total score and determine current phase
        total_score = 0
        subdomain_scores = {}
        
        # Extract subdomain scores from payload
        field_mapping = {
            "Governance_Q1_Numeric": ("Governance", "Visie op passende zorg"),
            "Governance_Q2_Numeric": ("Governance", "Leiderschap en eigenaarschap"),
            "Structuur_Q1_Numeric": ("Structuur", "Regionale samenwerking"),
            "Structuur_Q2_Numeric": ("Structuur", "Tools en platforms"),
            "Proces_Q1_Numeric": ("Proces", "Patiëntgericht procesontwerp"),
            "Proces_Q2_Numeric": ("Proces", "Leren en verbeteren"),
            "Uitkomsten_en_sturing_Q1_Numeric": ("Uitkomsten & sturing", "Outcomegericht werken"),
            "Uitkomsten_en_sturing_Q2_Numeric": ("Uitkomsten & sturing", "Monitoring en besluitvorming")
        }
        
        for field, (domain, subdomain) in field_mapping.items():
            if field in payload:
                score = int(payload[field])
                total_score += score
                if domain not in subdomain_scores:
                    subdomain_scores[domain] = []
                subdomain_scores[domain].append((subdomain, score))
        
        # Define maturity phases
        phases = [
            {"name": "Startfase", "range": "(0-14 pnt)", "emoji": "🌱", "description": "Eerste bewustwording, nog geen structuur", "focus": "Begrip en taal ontwikkelen", "color": "#FFCDD2"},
            {"name": "Aan de slag", "range": "(15-22 pnt)", "emoji": "🧱", "description": "Initiatieven starten, weinig samenhang", "focus": "Richting en partners bepalen", "color": "#FFE0B2"},
            {"name": "Op de kaart", "range": "(23-30 pnt)", "emoji": "🗺️", "description": "Structuur, rollen en data zijn aanwezig", "focus": "Governance en procesafspraken", "color": "#DCEDC8"},
            {"name": "In control", "range": "(31-36 pnt)", "emoji": "📊", "description": "Sturing op uitkomsten en processen", "focus": "PDCA-cyclus en datasturing", "color": "#C8E6C9"},
            {"name": "Voorloper", "range": "(37-40 pnt)", "emoji": "🚀", "description": "Passende zorg is geborgd en opgeschaald", "focus": "Leren, waardesturing, opschaling", "color": "#A5D6A7"}
        ]
        
        # Determine current phase
        current_phase = 0
        if total_score >= 37:
            current_phase = 4
        elif total_score >= 31:
            current_phase = 3
        elif total_score >= 23:
            current_phase = 2
        elif total_score >= 15:
            current_phase = 1
        else:
            current_phase = 0
        
        # Create the visualization - wider layout
        fig, ax = plt.subplots(figsize=(18, 6))  # Much wider, less tall
        ax.axis('off')
        
        # Define the layout: Left = Domain Table, Right = Score Bar Chart
        table_width = 0.4  # 40% for domain table
        chart_width = 0.55  # 55% for score chart
        gap = 0.05  # 5% gap between them
        
        # Domain colors matching client's scheme
        domain_colors = {
            'Governance': '#1565C0',      # Dark blue
            'Structuur': '#7B1FA2',       # Purple  
            'Proces': '#EF6C00',          # Orange
            'Uitkomsten & sturing': '#C62828'  # Red
        }
        
        # LEFT SIDE: Domain/Subdomain Table
        table_x = 0.02
        table_y_start = 0.85
        row_height = 0.08
        
        # Calculate total rows needed
        total_subdomains = sum(len(subs) for subs in subdomain_scores.values())
        table_height = len(subdomain_scores) * 0.12 + total_subdomains * 0.06
        
        current_y = table_y_start
        domain_order = ['Governance', 'Structuur', 'Proces', 'Uitkomsten & sturing']
        domain_icons = {
            'Governance': '🏛',
            'Structuur': '🧩', 
            'Proces': '🔧',
            'Uitkomsten & sturing': '📈'
        }
        
        # Draw domain table
        for domain in domain_order:
            if domain in subdomain_scores:
                # Domain header row
                rect = plt.Rectangle((table_x, current_y - 0.04), table_width, 0.08,
                                   facecolor=domain_colors[domain], alpha=0.9, 
                                   edgecolor='white', linewidth=2)
                ax.add_patch(rect)
                
                icon = domain_icons.get(domain, '⚙️')
                ax.text(table_x + 0.01, current_y, f"{icon} {domain}", 
                       ha='left', va='center', fontsize=12, fontweight='bold', color='white')
                
                current_y -= 0.08
                
                # Subdomain rows
                subdomains = subdomain_scores[domain]
                for subdomain, score in subdomains:
                    # Subdomain background (lighter version of domain color)
                    base_color = domain_colors[domain]
                    rect = plt.Rectangle((table_x, current_y - 0.03), table_width, 0.06,
                                       facecolor=base_color, alpha=0.3, 
                                       edgecolor='white', linewidth=1)
                    ax.add_patch(rect)
                    
                    # Subdomain text
                    ax.text(table_x + 0.02, current_y, f"• {subdomain}", 
                           ha='left', va='center', fontsize=10, color='#333333')
                    
                    # Score display
                    score_color = '#2E7D32' if score >= 4 else '#689F38' if score >= 3 else '#FF8F00' if score >= 2 else '#D32F2F'
                    ax.text(table_x + table_width - 0.02, current_y, f"{score}", 
                           ha='right', va='center', fontsize=11, fontweight='bold', 
                           color=score_color,
                           bbox=dict(boxstyle="circle,pad=0.2", facecolor='white', alpha=0.8))
                    
                    current_y -= 0.06
                
                current_y -= 0.02  # Extra space after each domain
        
        # RIGHT SIDE: Horizontal Score Progression Chart
        chart_x = table_x + table_width + gap
        chart_y_start = 0.85
        
        # Define score phases with exact client colors
        score_phases = [
            {"score_range": "1", "name": "Startfase", "range": "(0-14 pnt)", "color": "#FFCDD2", "description": "Eerste bewustwording, nog geen structuur"},
            {"score_range": "2", "name": "Aan de slag", "range": "(15-22 pnt)", "color": "#C8E6C9", "description": "Initiatieven starten, weinig samenhang"},
            {"score_range": "3", "name": "Op de kaart", "range": "(23-30 pnt)", "color": "#81C784", "description": "Structuur, rollen en data zijn aanwezig"},
            {"score_range": "4", "name": "In control", "range": "(31-36 pnt)", "color": "#4CAF50", "description": "Sturing op uitkomsten en processen"},
            {"score_range": "5", "name": "Voorloper", "range": "(37-40 pnt)", "color": "#2E7D32", "description": "Passende zorg is geborgd en opgeschaald"}
        ]
        
        # Determine current overall phase
        if total_score >= 37:
            current_phase_idx = 4
        elif total_score >= 31:
            current_phase_idx = 3
        elif total_score >= 23:
            current_phase_idx = 2
        elif total_score >= 15:
            current_phase_idx = 1
        else:
            current_phase_idx = 0
        
        # Draw horizontal progression bars
        bar_height = 0.12
        bar_y_start = 0.7
        
        for i, phase in enumerate(score_phases):
            y_pos = bar_y_start - (i * (bar_height + 0.02))
            
            # Highlight current phase
            is_current = i == current_phase_idx
            alpha = 1.0 if is_current else 0.5
            border_width = 3 if is_current else 1
            
            # Phase bar
            rect = plt.Rectangle((chart_x, y_pos), chart_width, bar_height,
                               facecolor=phase["color"], alpha=alpha, 
                               edgecolor='#333333', linewidth=border_width)
            ax.add_patch(rect)
            
            # Score number in circle
            ax.text(chart_x + 0.03, y_pos + bar_height/2, phase["score_range"], 
                   ha='center', va='center', fontsize=14, fontweight='bold',
                   bbox=dict(boxstyle="circle,pad=0.3", facecolor='white', alpha=0.9))
            
            # Phase name and range
            ax.text(chart_x + 0.08, y_pos + bar_height - 0.02, 
                   f"{phase['name']} {phase['range']}", 
                   ha='left', va='top', fontsize=11, fontweight='bold')
            
            # Description
            ax.text(chart_x + 0.08, y_pos + bar_height/2, 
                   phase['description'], 
                   ha='left', va='center', fontsize=9, color='#333333')
            
            # Add organization's subdomain scores within this phase
            matching_subdomains = []
            for domain, subdomains in subdomain_scores.items():
                for subdomain, score in subdomains:
                    if score == int(phase["score_range"]):
                        matching_subdomains.append(f"{subdomain}")
            
            if matching_subdomains:
                domains_text = " • ".join(matching_subdomains[:3])  # Limit to 3 for space
                if len(matching_subdomains) > 3:
                    domains_text += f" (+{len(matching_subdomains)-3} more)"
                ax.text(chart_x + 0.08, y_pos + 0.01, 
                       f"Uw organisatie: {domains_text}", 
                       ha='left', va='bottom', fontsize=8, style='italic', 
                       color='#1976D2', fontweight='bold')
        
        # Bottom status bar
        status_y = 0.08
        current_phase = score_phases[current_phase_idx]
        
        # Organization's current position
        status_rect = plt.Rectangle((chart_x, status_y), chart_width, 0.06,
                                  facecolor=current_phase["color"], alpha=0.9, 
                                  edgecolor='#333333', linewidth=2)
        ax.add_patch(status_rect)
        
        status_text = f"Uw positie: {current_phase['name']} (Totaal: {total_score} punten van 40)"
        ax.text(chart_x + chart_width/2, status_y + 0.03, status_text, 
               ha='center', va='center', fontsize=12, fontweight='bold', color='white')
        
        # Focus areas at the very bottom
        focus_areas = [
            "Begrip en taal ontwikkelen",
            "Richting en partners bepalen", 
            "Governance en procesafspraken",
            "PDCA-cyclus en datasturing",
            "Leren, waardesturing, opschaling"
        ]
        
        focus_y = 0.02
        focus_width = chart_width / len(focus_areas)
        
        for i, focus in enumerate(focus_areas):
            x_pos = chart_x + (i * focus_width)
            is_current_focus = i == current_phase_idx
            
            ax.text(x_pos + focus_width/2, focus_y, focus, 
                   ha='center', va='center', fontsize=7, 
                   fontweight='bold' if is_current_focus else 'normal',
                   rotation=15, color='#1976D2' if is_current_focus else '#666666')
        
        # Set limits
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        
        plt.tight_layout()
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                   facecolor='none', edgecolor='none', transparent=True, pad_inches=0.1)
        plt.close()
        
        logger.info(f"Score breakdown chart created: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logger.error(f"Error creating score breakdown chart: {str(e)}")
        raise

def create_domain_scores_table(domain_data, logger):
    """
    Create a table-style chart for domain scores
    """
    try:
        # Set up the figure with proper size
        fig, ax = plt.subplots(figsize=(12, 6))
        ax.axis('off')  # Hide axes
        
        # Define table dimensions
        rows = len(domain_data)
        cols = 3  # Index, Domain, Score, Rating
        
        # Define colors - back to original brighter colors
        header_color = '#4472C4'  # Blue header
        row_color_1 = '#F2F2F2'  # Light gray
        row_color_2 = '#FFFFFF'  # White
        
        # Create table data
        table_data = []
        
        # Add header
        headers = ['', 'Domains', 'Scores', 'Rating']
        
        # Add data rows with rating text
        for i, domain in enumerate(domain_data):
            row = [str(domain['index']), domain['name'], f"{domain['score']:.2f}", generate_star_rating(domain['score'])]
            table_data.append(row)
        
        # Calculate positions - wider table
        cell_height = 0.8
        cell_widths = [1.0, 3.5, 1.8, 3.5]  # Increased widths for better spacing
        total_width = sum(cell_widths)
        
        # Start position - center the table
        start_x = 0.2  # Reduced to center better
        start_y = len(table_data) * cell_height + 1
        
        # Draw header with original styling
        current_x = start_x
        for j, (header, width) in enumerate(zip(headers, cell_widths)):
            rect = Rectangle((current_x, start_y), width, cell_height, 
                           facecolor=header_color, edgecolor='white', linewidth=1)
            ax.add_patch(rect)
            
            # Add header text
            ax.text(current_x + width/2, start_y + cell_height/2, header,
                   ha='center', va='center', fontweight='bold', color='white', fontsize=11)
            current_x += width
        
        # Draw data rows
        for i, row in enumerate(table_data):
            current_x = start_x
            current_y = start_y - (i + 1) * cell_height
            
            # Determine row color - alternating colors for all rows
            row_color = row_color_1 if i % 2 == 0 else row_color_2
            text_color = 'black'
            font_weight = 'normal'
            
            # Draw cells with original styling
            for j, (cell_data, width) in enumerate(zip(row, cell_widths)):
                rect = Rectangle((current_x, current_y), width, cell_height,
                               facecolor=row_color, edgecolor='white', linewidth=1)
                ax.add_patch(rect)
                
                # Add cell text with original font sizing
                font_size = 10 if j == 3 else 11  # Smaller font for rating column
                ax.text(current_x + width/2, current_y + cell_height/2, cell_data,
                       ha='center', va='center', fontweight=font_weight, 
                       color=text_color, fontsize=font_size)
                current_x += width
        
        # Set the limits
        ax.set_xlim(0, total_width + 1)
        ax.set_ylim(0, start_y + cell_height + 0.5)
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name, dpi=300, bbox_inches='tight', 
                   facecolor='white', edgecolor='none')
        plt.close()
        
        logger.info(f"Domain scores table created and saved to: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logger.error(f"Error creating domain scores table: {str(e)}")
        raise

# Update the calculate_domain_scores function to work with the new data structure
def calculate_domain_scores(payload):
    """
    Calculate average scores for each domain (keeping for backward compatibility)
    """
    domains = {
        'Governance': [
            payload.get('Governance_Q1_Numeric', 0),
            payload.get('Governance_Q2_Numeric', 0)
        ],
        'Proces': [
            payload.get('Proces_Q1_Numeric', 0),
            payload.get('Proces_Q2_Numeric', 0)
        ],
        'Structuur': [
            payload.get('Structuur_Q1_Numeric', 0),
            payload.get('Structuur_Q2_Numeric', 0)
        ],
        'Uitkomsten & sturing': [
            payload.get('Uitkomsten_en_sturing_Q1_Numeric', 0),
            payload.get('Uitkomsten_en_sturing_Q2_Numeric', 0)
        ]
    }
    
    # Calculate averages
    domain_averages = {}
    for domain, scores in domains.items():
        scores = [int(score) for score in scores]
        domain_averages[domain] = sum(scores) / len(scores) if scores else 0
    
    return domain_averages

def get_transitiefase(total_sum):
    """
    Get transitiefase based on Total_Sum value
    """
    total_sum = int(total_sum)
    
    if 0 <= total_sum <= 14:
        return "Startfase"
    elif 15 <= total_sum <= 22:
        return "Aan de slag"
    elif 23 <= total_sum <= 30:
        return "Op de kaart"
    elif 31 <= total_sum <= 36:
        return "In control"
    elif 37 <= total_sum <= 40:
        return "Voorloper"
    else:
        return "Onbekend"

def get_lowest_scoring_domains(payload):
    """
    Get domains/questions with score 1 or 2 with keywords from headings
    """
    low_scoring = []
    
    # Keywords extraction mapping
    keywords = {
        'Governance_Q1_Numeric': "Visie",
        'Governance_Q2_Numeric': "Leiderschap",
        'Structuur_Q1_Numeric': "Samenwerking",
        'Structuur_Q2_Numeric': "Tools",
        'Proces_Q1_Numeric': "Patiëntgericht",
        'Proces_Q2_Numeric': "Leren",
        'Uitkomsten_en_sturing_Q1_Numeric': "Uitkomsten",
        'Uitkomsten_en_sturing_Q2_Numeric': "Data"
    }
    
    for key, keyword in keywords.items():
        score = int(payload.get(key, 0))
        if score in [1, 2]:
            low_scoring.append(f"{keyword}: {score}")
    
    return ", ".join(low_scoring) if low_scoring else "Geen lage scores"

def replace_placeholders(presentation, payload, logger):
    """
    Replace placeholders in PPT slides with lead data while preserving formatting
    """
    try:
        # Prepare replacement values
        organization = payload.get('Organization', '')
        first_name = payload.get('First_Name', '')
        last_name = payload.get('Last_Name', '')
        respondent_name = f"{first_name} {last_name}".strip()
        report_date = datetime.now().strftime('%d-%m-%Y')
        total_sum = payload.get('Total_Sum', '0')
        transitiefase = get_transitiefase(total_sum)
        lowest_domains = get_lowest_scoring_domains(payload)
        
        logger.info(f"Replacement values: Organization={organization}, Respondent={respondent_name}, Total_Sum={total_sum}, Transitiefase={transitiefase}")
        
        # Define replacements for each slide
        slide_replacements = {
            0: {  # Slide 1 (index 0)
                '{{organisatie}}': organization,
                '{{rapport_datum}}': report_date,
                '{{respondent_naam}}': respondent_name
            },
            3: {  # Slide 4 (index 3)
                '{{organisatie}}': organization,
                '{{totaalscore}}': str(total_sum),
                '{{transitiefase_naam}}': transitiefase
            },
            8: {  # Slide 9 (index 8)
                '{{organisatie}}': organization,
                '{{transitiefase}}': transitiefase,
                '{{laagst_scorende_domein}}': lowest_domains
            }
        }
        
        # Process each slide
        for slide_index, replacements in slide_replacements.items():
            if slide_index < len(presentation.slides):
                slide = presentation.slides[slide_index]
                logger.info(f"Processing slide {slide_index + 1}")
                
                # Replace text in all shapes on the slide
                for shape in slide.shapes:
                    # Check text in text frames for more complex shapes
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            # Get the full paragraph text first
                            paragraph_text = paragraph.text
                            
                            # Check if any placeholder exists in this paragraph
                            needs_replacement = False
                            new_paragraph_text = paragraph_text
                            
                            for placeholder, value in replacements.items():
                                if placeholder in paragraph_text:
                                    new_paragraph_text = new_paragraph_text.replace(placeholder, value)
                                    needs_replacement = True
                                    logger.info(f"Found {placeholder} in paragraph, replacing with {value}")
                            
                            # If replacement is needed, try to preserve formatting better
                            if needs_replacement:
                                # Clean the text to remove unwanted characters
                                cleaned_text = clean_text(new_paragraph_text)
                                
                                # Store all original run formatting
                                original_runs_formatting = []
                                for run in paragraph.runs:
                                    formatting = {
                                        'font_name': run.font.name,
                                        'font_size': run.font.size,
                                        'bold': run.font.bold,
                                        'italic': run.font.italic,
                                        'color': None
                                    }
                                    try:
                                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                                            formatting['color'] = run.font.color.rgb
                                    except:
                                        pass
                                    original_runs_formatting.append(formatting)
                                
                                # Clear all runs
                                paragraph.clear()
                                
                                # Add the new text with the best available formatting
                                new_run = paragraph.add_run()
                                new_run.text = cleaned_text
                                
                                # Apply the most common/first formatting
                                if original_runs_formatting:
                                    first_format = original_runs_formatting[0]
                                    try:
                                        if first_format['font_name']:
                                            new_run.font.name = first_format['font_name']
                                        if first_format['font_size']:
                                            new_run.font.size = first_format['font_size']
                                        if first_format['bold'] is not None:
                                            new_run.font.bold = first_format['bold']
                                        if first_format['italic'] is not None:
                                            new_run.font.italic = first_format['italic']
                                        if first_format['color']:
                                            new_run.font.color.rgb = first_format['color']
                                    except Exception as font_error:
                                        logger.warning(f"Could not apply formatting: {font_error}")
                    
                    # Fallback for simple text shapes
                    elif hasattr(shape, 'text'):
                        original_text = shape.text
                        new_text = original_text
                        
                        for placeholder, value in replacements.items():
                            if placeholder in new_text:
                                new_text = new_text.replace(placeholder, value)
                                logger.info(f"Replaced {placeholder} with {value} in simple text")
                        
                        if new_text != original_text:
                            # Clean the text to remove unwanted characters
                            cleaned_text = clean_text(new_text)
                            shape.text = cleaned_text
        
        logger.info("All placeholders processed successfully")
        return presentation
        
    except Exception as e:
        logger.error(f"Error replacing placeholders: {str(e)}")
        raise

def clean_text(text):
    """
    Clean unwanted characters from text while preserving intended line breaks and formatting
    """
    # Replace PowerPoint encoding characters with their actual intended characters
    replacements = {
        '_x000A': '\n',     # Line feed - replace with actual line break
        '_x000D': '\r',     # Carriage return - replace with carriage return
        '_x000B': '\n',     # Vertical tab - replace with line break
        '_x0009': '\t',     # Tab - replace with actual tab
        '\x0B': '\n',       # Actual vertical tab - replace with line break
        '\x0D\x0A': '\n',   # Windows line ending - replace with single line break
        '\r\n': '\n',       # Windows line ending - replace with single line break
    }
    
    cleaned = text
    for encoded_char, actual_char in replacements.items():
        cleaned = cleaned.replace(encoded_char, actual_char)
    
    # Clean up excessive line breaks (more than 2 consecutive)
    while '\n\n\n' in cleaned:
        cleaned = cleaned.replace('\n\n\n', '\n\n')
    
    # Clean up multiple spaces but preserve line breaks
    lines = cleaned.split('\n')
    cleaned_lines = []
    for line in lines:
        # Clean multiple spaces within each line
        while '  ' in line:
            line = line.replace('  ', ' ')
        cleaned_lines.append(line.strip())
    
    return '\n'.join(cleaned_lines)

# Run the Flask app when executed directly
if __name__ == "__main__":
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
